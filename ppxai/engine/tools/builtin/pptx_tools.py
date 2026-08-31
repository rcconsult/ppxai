"""
PowerPoint tools for multimodal attachments.

Phase 4.2 (v1.17.4). Two tools that let the model explore and read PPTX
files the user has attached via `/attach slides.pptx`:

    list_pptx_slides(file_id)
        → slide inventory with shape types (TEXT, TABLE, CHART, IMAGE)

    read_pptx_slide_text(file_id, slide)
        → text + tables from a specific slide as markdown

Both tools resolve `file_id` through the engine's SessionFileStore.
Guarded by `try: import pptx` (python-pptx package) at registration.

Slide rasterization (RenderPptxSlideTool via LibreOffice headless) and
embedded image extraction (ExtractPptxImagesTool) are deferred to a
follow-up step — they require system dependencies not all users have.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, List, Optional, Tuple

from ...types import ToolEngineProtocol, ToolManagerProtocol
from ..base import BaseTool
from ...file_ref import resolve_file_reference
from ...file_ref import FILE_REF_PROPERTIES
import base64
import subprocess
import tempfile
from ....common.libreoffice import find_libreoffice, libreoffice_available


_MAX_TEXT_CHARS = 100_000


def extract_pptx_slide_text(path: Path, slide_num: int) -> str:
    """Extract title + body text + tables from one PPTX slide as markdown.

    Pure, path-based public helper. Two callers today:
    - `ReadPptxSlideTextTool.execute()` — the LLM-facing tool, resolves
      `file_id` → `meta.path` then calls this
    - `server/routes/file_serve.py` — the path-based preview endpoint's
      LibreOffice-missing fallback (renders extracted text instead of
      raster slides). Lets the file-tree preview path degrade gracefully
      when LibreOffice is absent on the server.

    Returns the same markdown shape `ReadPptxSlideTextTool` returned
    historically (so the tool's downstream prompts don't shift). On
    error returns a short "Error: ..." string instead of raising —
    matches the existing tool's behavior. python-pptx is required.

    Args:
        path: Filesystem path to the .pptx/.ppt file.
        slide_num: 1-based slide index.

    Returns:
        Markdown string with `# {name} — Slide N of M`, optional
        `## Title`, body paragraphs, and `| ... |`-rendered tables.
        Returns `"Error: <reason>"` on any failure (python-pptx
        missing, file unreadable, slide_num out of range).
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed. Install with: pip install 'ppxai[data]'"

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return f"Error opening {path.name!r}: {exc}"

    total = len(prs.slides)
    try:
        slide_num = int(slide_num)
    except (TypeError, ValueError):
        return f"Error: slide must be an integer (got {slide_num!r})."

    if slide_num < 1 or slide_num > total:
        return f"Error: slide {slide_num} out of range (presentation has {total} slides)."

    target_slide = prs.slides[slide_num - 1]
    lines: List[str] = [f"# {path.name} — Slide {slide_num} of {total}\n"]

    # Title
    if target_slide.shapes.title and target_slide.shapes.title.has_text_frame:
        title = target_slide.shapes.title.text_frame.text.strip()
        if title:
            lines.append(f"## {title}\n")

    total_chars = sum(len(line) for line in lines)

    for shape in target_slide.shapes:
        if total_chars > _MAX_TEXT_CHARS:
            lines.append(
                f"\n[Output truncated at {_MAX_TEXT_CHARS:,} chars.]"
            )
            break

        if shape.has_table:
            table_md = _table_to_markdown(shape.table)
            lines.append(table_md + "\n")
            total_chars += len(table_md)

        elif shape.has_text_frame:
            # Skip the title shape (already rendered above)
            if (
                target_slide.shapes.title
                and shape.shape_id == target_slide.shapes.title.shape_id
            ):
                continue

            text = _text_frame_to_markdown(shape.text_frame)
            if text.strip():
                lines.append(text + "\n")
                total_chars += len(text)

        elif shape.has_chart:
            lines.append(f"[Chart: {shape.name}]\n")
            total_chars += 30

        elif hasattr(shape, "image") and shape.shape_type == 13:
            name = shape.name or "image"
            lines.append(f"[Image: {name}]\n")
            total_chars += 30

    result = "\n".join(lines)
    if not result.strip():
        return f"Slide {slide_num}: no text content found."

    return result


def _resolve_file(
    engine: Any,
    file_id: Optional[str] = None,
    path: Optional[str] = None,
) -> Tuple[Optional[Any], Optional[str]]:
    """Resolve a file reference via the unified engine resolver.

    Accepts EITHER `file_id` (SessionFileStore chat attachment) or
    `path` (workspace file). v1.18.7 — see `engine.file_ref` for the
    full resolution rules. Kept as a thin module-level shim so existing
    test fixtures using `SimpleNamespace(file_store=store)` keep
    working unchanged.
    """
    return resolve_file_reference(engine, file_id=file_id, path=path)


def _is_pptx(meta: Any) -> bool:
    """Check if a file is a PowerPoint presentation."""
    if "presentation" in (meta.media_type or ""):
        return True
    name = (meta.name or "").lower()
    return name.endswith((".pptx", ".ppt"))


def _shape_type_label(shape) -> str:
    """Human-readable label for a slide shape."""
    if shape.has_table:
        return "TABLE"
    if shape.has_chart:
        return "CHART"
    if hasattr(shape, "image") and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        return "IMAGE"
    if shape.has_text_frame:
        return "TEXT"
    return "OTHER"


class ListPptxSlidesTool(BaseTool):
    """List all slides in an attached PPTX file with shape inventories."""

    def __init__(self, engine: ToolEngineProtocol):
        self.engine = engine
        self.name = "list_pptx_slides"
        self.description = (
            "List all slides in a PowerPoint file. Returns slide numbers, "
            "titles (if any), and an inventory of shape types (TEXT, TABLE, "
            "CHART, IMAGE) per slide. Use this first to understand the "
            "presentation structure before reading specific slides. "
            "Pass either 'file_id' (for a presentation attached to this "
            "chat session, revived from the session cache on reload) or "
            "'path' (for a workspace file visible in the file tree, "
            "addressable from any session) — exactly one is required."
        )
        self.parameters = {
            "type": "object",
            "properties": dict(FILE_REF_PROPERTIES),
            "required": [],
        }

    async def execute(
        self,
        file_id: Optional[str] = None,
        path: Optional[str] = None,
        **kwargs,
    ) -> str:
        meta, err = _resolve_file(self.engine, file_id=file_id, path=path)
        if err:
            return f"Error: {err}"
        if not _is_pptx(meta):
            return f"Error: {meta.name!r} is not a PowerPoint file (type={meta.media_type!r})."

        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx not installed. Install with: pip install 'ppxai[data]'"

        try:
            prs = Presentation(str(meta.path))
        except Exception as exc:
            return f"Error opening {meta.name!r}: {exc}"

        slides = prs.slides
        total = len(slides)
        if total == 0:
            return f"{meta.name}: empty presentation (0 slides)."

        lines: List[str] = [f"# {meta.name} — {total} slide(s)\n"]

        for i, slide in enumerate(slides, 1):
            # Extract title from title placeholder if present
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()

            # Inventory shape types
            type_counts: dict = {}
            for shape in slide.shapes:
                label = _shape_type_label(shape)
                type_counts[label] = type_counts.get(label, 0) + 1

            shapes_str = ", ".join(
                f"{count}×{stype}" for stype, count in sorted(type_counts.items())
            ) or "(empty slide)"

            title_str = f" — {title}" if title else ""
            lines.append(
                f"## Slide {i}{title_str}\n"
                f"- Shapes: {shapes_str}\n"
            )

        return "\n".join(lines)


class ReadPptxSlideTextTool(BaseTool):
    """Read text + tables from a specific slide as markdown."""

    def __init__(self, engine: ToolEngineProtocol):
        self.engine = engine
        self.name = "read_pptx_slide_text"
        self.description = (
            "Read all text and tables from a specific slide of a PowerPoint "
            "file. Returns the content as markdown. Tables are rendered as "
            "markdown tables. Use list_pptx_slides first to see available "
            "slides. Pass either 'file_id' (chat attachment) or 'path' "
            "(workspace file) — exactly one is required."
        )
        self.parameters = {
            "type": "object",
            "properties": {
                **FILE_REF_PROPERTIES,
                "slide": {
                    "type": "integer",
                    "description": "1-indexed slide number.",
                },
            },
            "required": ["slide"],
        }

    async def execute(
        self,
        file_id: Optional[str] = None,
        path: Optional[str] = None,
        slide: int = 1,
        **kwargs,
    ) -> str:
        meta, err = _resolve_file(self.engine, file_id=file_id, path=path)
        if err:
            return f"Error: {err}"
        if not _is_pptx(meta):
            return f"Error: {meta.name!r} is not a PowerPoint file."
        # Delegate to the path-based public helper. The slide-name
        # prefix in the helper output uses `path.name` (== meta.name
        # for store-resolved files), so the historical output shape
        # is preserved byte-for-byte.
        return extract_pptx_slide_text(meta.path, slide)


def _text_frame_to_markdown(text_frame) -> str:
    """Convert a text frame's paragraphs to markdown text.

    Preserves paragraph breaks, bold/italic runs where possible,
    and bullet indicators from the paragraph level.
    """
    parts: List[str] = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Bullet detection: level > 0 or has bullet char
        level = para.level or 0
        if level > 0:
            indent = "  " * level
            parts.append(f"{indent}- {text}")
        else:
            parts.append(text)
    return "\n".join(parts)


def _table_to_markdown(table) -> str:
    """Convert a PPTX table to a markdown table."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip().replace("|", "\\|"))
        rows.append(cells)

    if not rows:
        return "(empty table)"

    header = rows[0]
    lines = []
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")

    for row in rows[1:]:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")

    return "\n".join(lines)


# =============================================================================
# Slide rasterization via LibreOffice headless
# =============================================================================


def _libreoffice_available() -> bool:
    """Check if LibreOffice headless is installed.

    Delegates to the shared cross-platform resolver — finds `libreoffice` or
    `soffice` on PATH, the macOS `.app` bundle, Windows Program Files, and the
    `PPXAI_LIBREOFFICE` override. (Previously `shutil.which("libreoffice")`
    only, which missed macOS's `soffice`.)
    """
    return libreoffice_available()


def render_pptx_slides(pptx_path: Path, cache_dir: Path) -> List[Path]:
    """Render all slides of a PPTX to PNG files via LibreOffice headless.

    Results are cached in *cache_dir* keyed by the source file name.
    The cache is invalidated when the source's mtime is newer than the
    oldest cached PNG — prevents the file-tree preview from showing
    stale slides after the model rewrites the .pptx (v1.18.7 fix).
    Returns a list of PNG paths sorted by slide number.
    """
    # Check cache first. v1.18.7: validate that the cache is at least as
    # new as the source — otherwise the file-tree preview keeps showing
    # the pre-edit slides after the model rewrites the .pptx.
    cached = sorted(cache_dir.glob("slide-*.png"))
    if cached:
        try:
            source_mtime = pptx_path.stat().st_mtime
        except OSError:
            # Source missing or unreadable. Keep the cache — at least the
            # user sees the last known render rather than nothing; a
            # re-render would just fail too.
            return cached
        try:
            oldest_cache_mtime = min(p.stat().st_mtime for p in cached)
        except OSError:
            # Cache PNG vanished between glob and stat — fall through
            # and re-render.
            oldest_cache_mtime = 0.0
        if oldest_cache_mtime >= source_mtime:
            return cached
        # Source is newer than the cache — invalidate so the next
        # render reflects the user's / model's edits.
        for p in cached:
            try:
                p.unlink()
            except OSError:
                pass

    cache_dir.mkdir(parents=True, exist_ok=True)

    soffice = find_libreoffice()
    if soffice is None:
        return []

    with tempfile.TemporaryDirectory() as tmpdir:
        # LibreOffice converts .pptx → .pdf headlessly; we then render
        # each PDF page to PNG with pypdfium2 (replaces the old
        # `pdftoppm` subprocess from poppler — pure-wheel, no system
        # binary required).
        # Single-step `--convert-to png` only emits the first slide on
        # some LibreOffice versions, so go via PDF.
        subprocess.run(
            [
                soffice, "--headless", "--norestore",
                "--convert-to", "pdf",
                "--outdir", tmpdir,
                str(pptx_path),
            ],
            capture_output=True,
            timeout=120,
        )
        # LibreOffice names output based on input filename
        pdf_candidates = list(Path(tmpdir).glob("*.pdf"))
        if not pdf_candidates:
            return []
        pdf_path = pdf_candidates[0]

        # v1.18.1: pypdfium2 replaces pdftoppm. Same 150 DPI as before;
        # `scale` is relative to 72 DPI (PDF's internal resolution).
        try:
            import pypdfium2 as pdfium  # noqa: PLC0415
        except ImportError:
            return []

        try:
            pdf = pdfium.PdfDocument(str(pdf_path))
        except pdfium.PdfiumError:
            return []

        try:
            for i in range(len(pdf)):
                # 1-indexed, zero-padded to 2 digits — preserves the
                # filename shape pdftoppm produced (slide-01.png ...).
                out = cache_dir / f"slide-{i + 1:02d}.png"
                page = pdf[i]
                bitmap = page.render(scale=150 / 72)
                try:
                    bitmap.to_pil().save(out, format="PNG")
                finally:
                    bitmap.close()
                    page.close()
        finally:
            pdf.close()

    return sorted(cache_dir.glob("slide-*.png"))


class SummarizePptxVisualTool(BaseTool):
    """Summarize a PPTX by rendering slides and captioning via VL model."""

    def __init__(self, engine: ToolEngineProtocol):
        self.engine = engine
        self.name = "summarize_pptx_visual"
        self.description = (
            "Visually summarize a PowerPoint file by rendering all slides "
            "to images and sending them to a vision-language model. Returns "
            "a description of each slide's visual content. Much faster than "
            "reading slides one-by-one with read_pptx_slide_text. Requires "
            "LibreOffice and a vision model. Pass either 'file_id' (chat "
            "attachment) or 'path' (workspace file) — exactly one is required."
        )
        self.parameters = {
            "type": "object",
            "properties": {
                **FILE_REF_PROPERTIES,
                "slides": {
                    "type": "string",
                    "description": "Slide range: 'all' or '1-5' or '3,7,12'. Default: all.",
                    "default": "all",
                },
            },
            "required": [],
        }

    async def execute(
        self,
        file_id: Optional[str] = None,
        path: Optional[str] = None,
        slides: str = "all",
        **kwargs,
    ) -> str:
        meta, err = _resolve_file(self.engine, file_id=file_id, path=path)
        if err:
            return f"Error: {err}"
        if not _is_pptx(meta):
            return f"Error: {meta.name!r} is not a PowerPoint file."

        # Render slides
        cache_dir = meta.path.parent / "slides"
        try:
            pngs = render_pptx_slides(meta.path, cache_dir)
        except Exception as exc:
            return f"Error rendering slides: {exc}"

        if not pngs:
            return "Error: no slides rendered."

        # Parse slide selection
        total = len(pngs)
        if slides == "all":
            indices = list(range(total))
        else:
            indices = []
            for part in slides.replace(" ", "").split(","):
                if "-" in part:
                    a, b = part.split("-", 1)
                    indices.extend(range(int(a) - 1, min(int(b), total)))
                else:
                    idx = int(part) - 1
                    if 0 <= idx < total:
                        indices.append(idx)

        if not indices:
            return f"Error: no valid slides in range '{slides}' (total: {total})."

        # Caption each slide via VL sidecar
        if not hasattr(self.engine, "caption_image"):
            return "Error: engine does not support caption_image."

        results: List[str] = [f"# {meta.name} — Visual Summary ({len(indices)} slides)\n"]

        for idx in indices:
            slide_num = idx + 1
            png_bytes = pngs[idx].read_bytes()
            caption = self.engine.caption_image(
                f"{meta.name} slide {slide_num}",
                "image/png",
                png_bytes,
            )
            if not caption:
                caption = "(VL model returned no description)"
            results.append(f"## Slide {slide_num}\n{caption}\n")

        return "\n".join(results)


class RenderPptxSlideTool(BaseTool):
    """Render a PPTX slide to a PNG image via LibreOffice headless."""

    def __init__(self, engine: ToolEngineProtocol):
        self.engine = engine
        self.name = "render_pptx_slide"
        self.description = (
            "Render a specific slide from a PowerPoint file as a PNG image. "
            "Returns the image as a base64 data URI. Requires LibreOffice "
            "headless (installed in the container). Use list_pptx_slides "
            "first to see available slides. Pass either 'file_id' (chat "
            "attachment) or 'path' (workspace file) — exactly one is required."
        )
        self.parameters = {
            "type": "object",
            "properties": {
                **FILE_REF_PROPERTIES,
                "slide": {
                    "type": "integer",
                    "description": "Slide number (1-based).",
                    "default": 1,
                },
            },
            "required": [],
        }

    async def execute(
        self,
        file_id: Optional[str] = None,
        path: Optional[str] = None,
        slide: int = 1,
        **kwargs,
    ) -> str:
        if not _libreoffice_available():
            return "Error: LibreOffice is not installed. Slide rendering unavailable."

        meta, err = _resolve_file(self.engine, file_id=file_id, path=path)
        if err:
            return f"Error: {err}"
        if not _is_pptx(meta):
            return f"Error: {meta.name!r} is not a PowerPoint file."

        cache_dir = meta.path.parent / "slides"
        try:
            pngs = render_pptx_slides(meta.path, cache_dir)
        except subprocess.TimeoutExpired:
            return "Error: LibreOffice timed out rendering slides."
        except Exception as exc:
            return f"Error rendering slides: {exc}"

        if not pngs:
            return "Error: LibreOffice produced no output. The file may be corrupt."

        if slide < 1 or slide > len(pngs):
            return f"Error: slide {slide} out of range (1-{len(pngs)})."

        png_path = pngs[slide - 1]
        png_bytes = png_path.read_bytes()

        # v1.18.7: tool-produced multimodal artifacts go through the
        # SessionFileStore — same lifecycle as user-uploaded chat
        # attachments (revived from cache on session reload). Returning
        # the inline base64 data URI used to cost ~80-110K tokens per
        # call (rasterized slides at 150 DPI), enough to blow the
        # context window on a single iteration. Now we save the PNG
        # and return a compact reference; visual analysis is a
        # follow-up tool call.
        file_store = getattr(self.engine, "file_store", None)
        if file_store is None:
            # No store available (some test fixtures, unconfigured
            # engines) — fall back to inline data URI for compat.
            b64 = base64.b64encode(png_bytes).decode()
            return f"data:image/png;base64,{b64}"

        artifact_name = f"{Path(meta.name).stem}_slide_{slide}.png"
        try:
            saved = file_store.save(
                artifact_name, png_bytes, media_type="image/png"
            )
        except OSError as exc:
            return f"Error: failed to save rendered slide: {exc}"

        size_kb = len(png_bytes) / 1024
        return (
            f"Rendered slide {slide} of {meta.name} "
            f"({size_kb:.1f} KB PNG).\n"
            f"Saved to session attachments: file_id={saved.file_id}\n\n"
            f"The PNG is now in the session file store (same lifecycle "
            f"as chat-uploaded files, revived on session reload). To "
            f"analyze it visually, use summarize_pptx_visual on the "
            f"original .pptx (file_id or path), which captions all "
            f"slides via the vision sidecar. The plain shape inventory "
            f"and text content are available without re-rendering via "
            f"list_pptx_slides and read_pptx_slide_text."
        )


# =============================================================================
# Registration
# =============================================================================


def register_tools(manager: ToolManagerProtocol, engine: ToolEngineProtocol) -> bool:
    """Register PPTX tools. Returns False if python-pptx is not installed."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        return False

    if engine is None:
        return False

    manager.register_tool(ListPptxSlidesTool(engine))
    manager.register_tool(ReadPptxSlideTextTool(engine))
    if _libreoffice_available():
        manager.register_tool(RenderPptxSlideTool(engine))
        # VL visual summary — requires both LibreOffice and a VL sidecar
        if hasattr(engine, "has_vision_sidecar") and engine.has_vision_sidecar():
            manager.register_tool(SummarizePptxVisualTool(engine))
    return True


__all__ = [
    "ListPptxSlidesTool",
    "ReadPptxSlideTextTool",
    "RenderPptxSlideTool",
    "SummarizePptxVisualTool",
    "render_pptx_slides",
    "register_tools",
]
