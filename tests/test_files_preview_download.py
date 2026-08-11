"""HTTP route tests for /files/preview (path-based) + /files/download.

v1.18.7. Covers the two new path-based endpoints filed alongside the
office-doc preview wiring (regression: file-tree click on PPTX/DOCX
hit `CodeEditorView` and returned "Cannot read binary file"; structural
fix wires them through a real preview path).

Three concerns covered:

1. **/files/preview?path= happy paths** — spreadsheet rejection (goes
   through /files/read instead), PPTX/Word metadata via ?total=true,
   PPTX/Word LibreOffice-fallback to text_fallback JSON.

2. **/files/download?path= happy + error paths** — Content-Disposition:
   attachment, MIME type lookup, 404 on missing, 403 on out-of-tree,
   400 on directory, 400 on empty path.

3. **/files/read office-spreadsheet branch** — xlsx/xls/csv now return
   base64+mime_type instead of falling through to text-decode 400.

The LibreOffice-fallback test branches are runtime-dispatched: if
LibreOffice is installed on the test host the rendering path is
exercised; if not, the text-extraction path. Both end states are
asserted so the test runs green either way.
"""

from __future__ import annotations

import base64
import json
from pathlib import Path

import pytest

pytest.importorskip("fastapi")
pytest.importorskip("httpx")

from fastapi.testclient import TestClient


@pytest.fixture
def http_client():
    import ppxai.server.http as http_module
    with TestClient(http_module.app, raise_server_exceptions=False) as client:
        yield client


def _session(name: str) -> dict:
    return {"X-Session-Id": f"preview-dl-{name}"}


def _anchor_to(client: TestClient, headers: dict, path: Path) -> None:
    client.post(
        "/context/working_dir",
        json={"path": str(path)},
        headers=headers,
    )


def _make_minimal_xlsx(tmp_path: Path) -> Path:
    """Tiny valid XLSX file. Uses openpyxl if available, else a
    pre-baked byte blob (smallest possible OOXML structure)."""
    target = tmp_path / "sample.xlsx"
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "header"
    ws["A2"] = "value"
    wb.save(str(target))
    return target


# ---------------------------------------------------------------------------
# /files/read — office-spreadsheet branch (v1.18.7 extension)
# ---------------------------------------------------------------------------


class TestReadOfficeSpreadsheet:
    """xlsx/xls/csv used to fall through to text-decode 400.
    v1.18.7 added them to BINARY_PREVIEW_EXTENSIONS so /files/read
    returns base64 + mime_type, same shape as image/pdf."""

    def test_xlsx_returns_base64_and_office_spreadsheet_type(self, http_client, tmp_path):
        target = _make_minimal_xlsx(tmp_path)
        headers = _session("xlsx-base64")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.post(
            "/files/read",
            json={"path": "sample.xlsx", "cwd_anchor": str(tmp_path)},
            headers=headers,
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["type"] == "office_spreadsheet"
        assert body["mime_type"] == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        assert body["filename"] == "sample.xlsx"
        # content is base64 of the file bytes
        decoded = base64.b64decode(body["content"])
        assert decoded == target.read_bytes()

    def test_csv_returns_base64_and_office_spreadsheet_type(self, http_client, tmp_path):
        target = tmp_path / "data.csv"
        target.write_text("a,b,c\n1,2,3\n", encoding="utf-8")
        headers = _session("csv-base64")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.post(
            "/files/read",
            json={"path": "data.csv"},
            headers=headers,
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["type"] == "office_spreadsheet"
        assert body["mime_type"] == "text/csv"
        # CSV bytes round-trip via base64
        decoded = base64.b64decode(body["content"])
        # On Windows the write may have CRLF — compare bytes directly
        assert decoded == target.read_bytes()


# ---------------------------------------------------------------------------
# /files/read — PPTX/DOCX error message points at /files/preview
# ---------------------------------------------------------------------------


class TestReadOfficePresentationHint:
    """PPTX + DOCX intentionally aren't in BINARY_PREVIEW_EXTENSIONS
    (they need LibreOffice conversion, handled by /files/preview).
    /files/read returns 400 with a HINT pointing at the right endpoint
    so client developers don't have to guess."""

    def test_pptx_400_mentions_preview_endpoint(self, http_client, tmp_path):
        # Don't need a real PPTX — the route opens the file via read_text
        # and fails on UnicodeDecodeError before reading the bytes.
        # Use a byte sequence guaranteed to fail UTF-8 decode.
        target = tmp_path / "fake.pptx"
        target.write_bytes(b"\xff\xfe\xfd")  # invalid UTF-8 start bytes
        headers = _session("pptx-hint")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.post(
            "/files/read",
            json={"path": "fake.pptx"},
            headers=headers,
        )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert "/files/preview?path=" in detail
        assert ".pptx" in detail

    def test_docx_400_mentions_preview_endpoint(self, http_client, tmp_path):
        target = tmp_path / "fake.docx"
        target.write_bytes(b"\xff\xfe\xfd")
        headers = _session("docx-hint")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.post(
            "/files/read",
            json={"path": "fake.docx"},
            headers=headers,
        )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert "/files/preview?path=" in detail


# ---------------------------------------------------------------------------
# /files/preview?path= — input validation
# ---------------------------------------------------------------------------


class TestPreviewByPathValidation:
    """Unsupported extensions are rejected with a 400 that names the
    valid set; missing path is rejected at FastAPI Query level (422)."""

    def test_txt_rejected_as_unsupported(self, http_client, tmp_path):
        target = tmp_path / "note.txt"
        target.write_text("hello", encoding="utf-8")
        headers = _session("preview-txt")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/preview?path=note.txt&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert "Unsupported extension" in detail
        # Hint about where text/image/pdf goes instead
        assert "/files/read" in detail

    def test_xlsx_rejected_as_unsupported_for_preview(self, http_client, tmp_path):
        """Spreadsheets go through /files/read (SheetJS client-side),
        not /files/preview. The error message should make this clear."""
        target = _make_minimal_xlsx(tmp_path)
        headers = _session("preview-xlsx-reject")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/preview?path=sample.xlsx&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 400
        assert "Unsupported extension" in resp.json()["detail"]

    def test_missing_file_returns_404(self, http_client, tmp_path):
        headers = _session("preview-missing")
        _anchor_to(http_client, headers, tmp_path)
        resp = http_client.get(
            f"/files/preview?path=nope.pptx&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 404


# ---------------------------------------------------------------------------
# /files/download?path=
# ---------------------------------------------------------------------------


class TestDownload:
    """Path-based file download with Content-Disposition: attachment.
    Reuses _resolve_safe_path so security is identical to /files/read."""

    def test_text_file_returns_bytes_with_attachment_disposition(self, http_client, tmp_path):
        target = tmp_path / "report.txt"
        target.write_text("hello world", encoding="utf-8")
        headers = _session("dl-text")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/download?path=report.txt&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 200, resp.text
        assert resp.content == target.read_bytes()
        cd = resp.headers.get("content-disposition", "")
        assert "attachment" in cd
        assert 'filename="report.txt"' in cd

    def test_xlsx_uses_office_mime_type(self, http_client, tmp_path):
        target = _make_minimal_xlsx(tmp_path)
        headers = _session("dl-xlsx")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/download?path=sample.xlsx&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 200
        assert resp.headers["content-type"] == (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        assert resp.content == target.read_bytes()

    def test_unknown_extension_uses_octet_stream(self, http_client, tmp_path):
        target = tmp_path / "blob.weird"
        target.write_bytes(b"\x00\x01\x02")
        headers = _session("dl-weird")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/download?path=blob.weird&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 200
        assert resp.headers["content-type"] == "application/octet-stream"

    def test_missing_file_returns_404(self, http_client, tmp_path):
        headers = _session("dl-missing")
        _anchor_to(http_client, headers, tmp_path)
        resp = http_client.get(
            f"/files/download?path=does-not-exist.txt&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 404

    def test_directory_path_returns_400(self, http_client, tmp_path):
        sub = tmp_path / "src"
        sub.mkdir()
        headers = _session("dl-dir")
        _anchor_to(http_client, headers, tmp_path)
        resp = http_client.get(
            f"/files/download?path=src&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 400

    def test_empty_path_returns_400(self, http_client, tmp_path):
        headers = _session("dl-empty")
        _anchor_to(http_client, headers, tmp_path)
        resp = http_client.get(
            "/files/download?path=",
            headers=headers,
        )
        assert resp.status_code == 400


# ---------------------------------------------------------------------------
# /files/preview?path= — LibreOffice-or-fallback dispatch
# ---------------------------------------------------------------------------


def _libreoffice_present() -> bool:
    """Reuse the server's own probe so the test branches the same way
    the route does (avoids host-skew where the test thinks LO is missing
    but the route finds it, or vice versa)."""
    from ppxai.engine.tools.builtin.pptx_tools import _libreoffice_available
    return _libreoffice_available()


def _libreoffice_can_render(tmp_path) -> bool:
    """True iff LibreOffice can actually rasterize a file under *tmp_path*.

    `_libreoffice_present()` only proves the binary exists — a snap-confined
    install can't read `/tmp` and produces no output (exit 0), so the route
    degrades to text_fallback exactly as if LibreOffice were absent. Branch the
    render assertions on real capability so a confined dev box takes the same
    text_fallback branch the route does, instead of a false failure. The coder
    image's apt libreoffice is unconfined and renders (verified in-pod)."""
    from ppxai.common.libreoffice import libreoffice_can_read
    return _libreoffice_present() and libreoffice_can_read(tmp_path)


class TestPreviewByPathPptx:
    """LibreOffice and text-fallback both end with a usable response;
    we assert the shape the client expects in each case."""

    def _make_minimal_pptx(self, tmp_path: Path) -> Path:
        pptx_mod = pytest.importorskip("pptx")
        prs = pptx_mod.Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        # Add a text box so extracted text has content to find
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "Hello PPTX"
        target = tmp_path / "sample.pptx"
        prs.save(str(target))
        return target

    def test_total_query_returns_metadata(self, http_client, tmp_path):
        target = self._make_minimal_pptx(tmp_path)
        headers = _session("pptx-total")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/preview?path=sample.pptx&total=true&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["name"] == "sample.pptx"
        assert body["total"] == 1
        assert body["kind"] == "presentation"
        # When LibreOffice is present: type=pptx; when absent: still pptx
        # but with libreoffice_available=False.
        # When LibreOffice can actually render: type=pptx + count via raster.
        # When it can't (absent, or present-but-confined): the route degrades to
        # a text_fallback whose libreoffice_available is False. Either way the
        # response is a usable 200 (never a 500) — that degrade is the contract.
        if _libreoffice_can_render(tmp_path):
            assert body["type"] == "pptx"
        else:
            assert body["type"] in ("pptx", "text_fallback")
            assert body.get("libreoffice_available") is False

    def test_slide_fetch_returns_png_or_text_fallback(self, http_client, tmp_path):
        target = self._make_minimal_pptx(tmp_path)
        headers = _session("pptx-slide")
        _anchor_to(http_client, headers, tmp_path)

        resp = http_client.get(
            f"/files/preview?path=sample.pptx&slide=1&cwd_anchor={tmp_path}",
            headers=headers,
        )
        assert resp.status_code == 200, resp.text
        if _libreoffice_can_render(tmp_path):
            assert resp.headers["content-type"] == "image/png"
            # PNG magic bytes
            assert resp.content[:8] == b"\x89PNG\r\n\x1a\n"
        else:
            # LibreOffice absent, OR present-but-unable-to-render (snap
            # confinement): the route degrades to text_fallback either way.
            body = resp.json()
            assert body["type"] == "text_fallback"
            assert body["kind"] == "presentation"
            assert body["slide"] == 1
            assert body["libreoffice_available"] is False
            # extract_pptx_slide_text returns "# <name> — Slide 1 of N" header
            assert "Slide 1 of 1" in body["content"]


class TestUnifiedPreviewContract:
    """Item 26 (v1.18.8): both /files/preview routes delegate to the shared
    `render_office_preview` helper, which has ONE shape and NEVER 503/500s for
    missing LibreOffice — legacy `.ppt`/`.doc` included (they used to hit a
    python-pptx/docx 500 or a hard 503). These exercise the helper directly so
    they run without office libs or a file_store.
    """

    def _no_libreoffice(self, monkeypatch):
        # The helper does `from ...pptx_tools import _libreoffice_available`
        # at call time, so patching the module attribute is picked up.
        monkeypatch.setattr(
            "ppxai.engine.tools.builtin.pptx_tools._libreoffice_available",
            lambda: False,
        )

    def test_legacy_ppt_without_libreoffice_is_text_fallback_not_503(self, monkeypatch, tmp_path):
        from ppxai.server.routes.files import render_office_preview
        self._no_libreoffice(monkeypatch)
        resp = render_office_preview(tmp_path / "deck.ppt", "deck.ppt", ".ppt", tmp_path, total=True)
        assert resp.status_code == 200
        body = json.loads(resp.body)
        assert body == {
            "type": "text_fallback",
            "kind": "presentation",
            "content": body["content"],  # message text, asserted below
            "name": "deck.ppt",
            "total": 1,
            "libreoffice_available": False,
        }
        assert "LibreOffice" in body["content"]

    def test_legacy_doc_without_libreoffice_is_text_fallback_word(self, monkeypatch, tmp_path):
        from ppxai.server.routes.files import render_office_preview
        self._no_libreoffice(monkeypatch)
        resp = render_office_preview(tmp_path / "memo.doc", "memo.doc", ".doc", tmp_path, total=True)
        assert resp.status_code == 200
        body = json.loads(resp.body)
        assert body["type"] == "text_fallback"
        assert body["kind"] == "word"
        assert body["libreoffice_available"] is False

    def test_total_response_always_has_unified_keys(self, monkeypatch, tmp_path):
        # Every total/metadata response carries the full key set regardless of
        # branch — the contract the clients now rely on.
        from ppxai.server.routes.files import render_office_preview
        self._no_libreoffice(monkeypatch)
        resp = render_office_preview(tmp_path / "x.ppt", "x.ppt", ".ppt", tmp_path, total=True)
        body = json.loads(resp.body)
        for key in ("type", "kind", "name", "total", "libreoffice_available"):
            assert key in body, f"unified preview shape missing {key!r}"

    def _make_real_pptx(self, tmp_path):
        pptx_mod = pytest.importorskip("pptx")
        from pptx.util import Inches
        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Hello PPTX"
        target = tmp_path / "deck.pptx"
        prs.save(str(target))
        return target

    def test_pptx_present_libreoffice_empty_render_degrades_not_500(
        self, monkeypatch, tmp_path
    ):
        # A present-but-confined/broken LibreOffice renders NOTHING (the Ubuntu
        # snap can't read /tmp and exits 0 with no output). The route used to
        # raise a hard 500 "No slides rendered", breaking the "never 503/500 for
        # a preview we can't rasterize" contract. It must instead degrade to the
        # SAME text_fallback the LibreOffice-missing path uses. (The coder image
        # ships an unconfined apt libreoffice and renders a real PNG — verified
        # in-pod 2026-08-11; this fences the degrade for confined/broken hosts.)
        from ppxai.server.routes.files import render_office_preview
        target = self._make_real_pptx(tmp_path)
        monkeypatch.setattr(
            "ppxai.engine.tools.builtin.pptx_tools._libreoffice_available",
            lambda: True,
        )
        monkeypatch.setattr(
            "ppxai.engine.tools.builtin.pptx_tools.render_pptx_slides",
            lambda *a, **k: [],  # present LibreOffice, empty render
        )
        # slide fetch: text_fallback, not 500
        resp = render_office_preview(target, "deck.pptx", ".pptx", tmp_path, slide=1)
        assert resp.status_code == 200
        body = json.loads(resp.body)
        assert body["type"] == "text_fallback"
        assert body["kind"] == "presentation"
        assert body["libreoffice_available"] is False
        assert "Slide 1 of 1" in body["content"]
        # total/metadata: also 200, never 500
        resp_total = render_office_preview(
            target, "deck.pptx", ".pptx", tmp_path, total=True
        )
        assert resp_total.status_code == 200
        body_total = json.loads(resp_total.body)
        assert body_total["type"] == "pptx"
        assert body_total["total"] == 1
        assert body_total["libreoffice_available"] is False

    def test_unsupported_extension_rejected_400(self, tmp_path):
        from ppxai.server.routes.files import render_office_preview
        from fastapi import HTTPException
        with pytest.raises(HTTPException) as exc:
            render_office_preview(tmp_path / "x.txt", "x.txt", ".txt", tmp_path)
        assert exc.value.status_code == 400
