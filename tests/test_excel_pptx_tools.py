"""Tests for Excel and PowerPoint tools (Phase 4, v1.17.4).

Creates real .xlsx and .pptx files in tmp_path via openpyxl and
python-pptx, registers them with SessionFileStore, and invokes the
tools through their execute() coroutines — same test pattern as
test_pdf_tools.py.
"""

from __future__ import annotations

import asyncio
import io
from types import SimpleNamespace

import pytest

from ppxai.engine.session_store import SessionFileStore

openpyxl = pytest.importorskip("openpyxl")
pptx_mod = pytest.importorskip("pptx")

from ppxai.engine.tools.builtin.excel_tools import (  # noqa: E402
    ListExcelSheetsTool,
    ReadExcelSheetTool,
)
from ppxai.engine.tools.builtin.excel_tools import (
    register_tools as register_excel,
)
from ppxai.engine.tools.builtin.pptx_tools import (  # noqa: E402
    ListPptxSlidesTool,
    ReadPptxSlideTextTool,
)
from ppxai.engine.tools.builtin.pptx_tools import (
    register_tools as register_pptx,
)

# -----------------------------------------------------------------------------
# Test data builders
# -----------------------------------------------------------------------------


def _make_xlsx(sheets: dict[str, list[list]]) -> bytes:
    """Build a real .xlsx from {sheet_name: [[row], ...]} data."""
    wb = openpyxl.Workbook()
    # Remove the default sheet
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slides: list[dict]) -> bytes:
    """Build a real .pptx from a list of slide dicts.

    Each dict can have:
        title: str — slide title
        content: str — body text
        table: list[list[str]] — a table (first row = header)
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Title
        title = slide_data.get("title")
        if title:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = title

        # Body content
        content = slide_data.get("content")
        if content:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.text = content

        # Table
        table_data = slide_data.get("table")
        if table_data and len(table_data) >= 1:
            rows = len(table_data)
            cols = len(table_data[0])
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(3), Inches(8), Inches(2)
            )
            tbl = table_shape.table
            for r, row in enumerate(table_data):
                for c, cell_text in enumerate(row):
                    tbl.cell(r, c).text = str(cell_text)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def store(tmp_path) -> SessionFileStore:
    return SessionFileStore(base_dir=tmp_path / "uploads")


@pytest.fixture
def fake_engine(store):
    return SimpleNamespace(file_store=store)


@pytest.fixture
def sample_xlsx_meta(store):
    xlsx = _make_xlsx({
        "Sales": [
            ["Region", "Q1", "Q2", "Q3", "Q4"],
            ["North", 100, 150, 130, 180],
            ["South", 80, 90, 110, 120],
            ["East", 200, 210, 195, 250],
        ],
        "Summary": [
            ["Total", 1815],
        ],
    })
    return store.save(
        "report.xlsx", xlsx,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


@pytest.fixture
def sample_pptx_meta(store):
    pptx_bytes = _make_pptx([
        {"title": "Introduction", "content": "Welcome to the presentation"},
        {
            "title": "Data Overview",
            "content": "Key metrics for Q1",
            "table": [
                ["Metric", "Value"],
                ["Revenue", "$1.2M"],
                ["Users", "50K"],
            ],
        },
        {"title": "Conclusion", "content": "Thank you for your time"},
    ])
    return store.save(
        "slides.pptx", pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# -----------------------------------------------------------------------------
# Excel: ListExcelSheetsTool
# -----------------------------------------------------------------------------


class TestListExcelSheets:
    def test_lists_all_sheets(self, fake_engine, sample_xlsx_meta):
        tool = ListExcelSheetsTool(fake_engine)
        result = asyncio.run(tool.execute(file_id=sample_xlsx_meta.file_id))
        assert "report.xlsx" in result
        assert "2 sheet(s)" in result
        assert "Sales" in result
        assert "Summary" in result
        # Dimensions visible
        assert "Rows:" in result
        assert "Columns:" in result
        # Headers visible
        assert "Region" in result

    def test_missing_file_id(self, fake_engine):
        tool = ListExcelSheetsTool(fake_engine)
        result = asyncio.run(tool.execute(file_id="nonexistent"))
        assert "Error" in result
        assert "Unknown file_id" in result

    def test_non_excel_rejected(self, store):
        meta = store.save("image.png", b"\x89PNG\r\n\x1a\n" + b"\x00" * 50,
                          media_type="image/png")
        tool = ListExcelSheetsTool(SimpleNamespace(file_store=store))
        result = asyncio.run(tool.execute(file_id=meta.file_id))
        assert "Error" in result
        assert "not an Excel file" in result


# -----------------------------------------------------------------------------
# Excel: ReadExcelSheetTool
# -----------------------------------------------------------------------------


class TestReadExcelSheet:
    def test_reads_as_markdown(self, fake_engine, sample_xlsx_meta):
        tool = ReadExcelSheetTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_xlsx_meta.file_id, sheet="Sales"
        ))
        assert "report.xlsx" in result
        assert "Sales" in result
        # Markdown table format
        assert "| Region" in result
        assert "| ---" in result
        assert "| North" in result
        assert "| 100" in result or "| 100 " in result

    def test_reads_as_csv(self, fake_engine, sample_xlsx_meta):
        tool = ReadExcelSheetTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_xlsx_meta.file_id, sheet="Sales", as_markdown=False
        ))
        # CSV format — comma-separated
        assert "Region," in result
        assert "North," in result

    def test_sheet_not_found(self, fake_engine, sample_xlsx_meta):
        tool = ReadExcelSheetTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_xlsx_meta.file_id, sheet="Nonexistent"
        ))
        assert "Error" in result
        assert "not found" in result
        # Lists available sheets
        assert "Sales" in result

    def test_row_limit(self, fake_engine, store):
        # Large sheet (20 rows), read with limit of 5
        rows = [["Col1", "Col2"]] + [[f"row{i}", i] for i in range(20)]
        xlsx = _make_xlsx({"Big": rows})
        meta = store.save(
            "big.xlsx", xlsx,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        tool = ReadExcelSheetTool(SimpleNamespace(file_store=store))
        result = asyncio.run(tool.execute(
            file_id=meta.file_id, sheet="Big", rows=5
        ))
        # Should truncate and show a hint
        assert "Showing" in result or "rows" in result.lower()
        # First few rows present
        assert "row0" in result
        assert "row4" in result
        # Later rows absent
        assert "row15" not in result


# -----------------------------------------------------------------------------
# PPTX: ListPptxSlidesTool
# -----------------------------------------------------------------------------


class TestListPptxSlides:
    def test_lists_all_slides(self, fake_engine, sample_pptx_meta):
        tool = ListPptxSlidesTool(fake_engine)
        result = asyncio.run(tool.execute(file_id=sample_pptx_meta.file_id))
        assert "slides.pptx" in result
        assert "3 slide(s)" in result
        assert "Slide 1" in result
        assert "Slide 2" in result
        assert "Slide 3" in result
        # Shape types visible
        assert "TEXT" in result
        # Slide 2 has a table
        assert "TABLE" in result

    def test_missing_file_id(self, fake_engine):
        tool = ListPptxSlidesTool(fake_engine)
        result = asyncio.run(tool.execute(file_id="ghost"))
        assert "Error" in result

    def test_non_pptx_rejected(self, store):
        meta = store.save("doc.pdf", b"%PDF-1.4\n" + b"\x00" * 50,
                          media_type="application/pdf")
        tool = ListPptxSlidesTool(SimpleNamespace(file_store=store))
        result = asyncio.run(tool.execute(file_id=meta.file_id))
        assert "Error" in result
        assert "not a PowerPoint" in result


# -----------------------------------------------------------------------------
# PPTX: ReadPptxSlideTextTool
# -----------------------------------------------------------------------------


class TestReadPptxSlideText:
    def test_reads_text_from_slide(self, fake_engine, sample_pptx_meta):
        tool = ReadPptxSlideTextTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_pptx_meta.file_id, slide=1
        ))
        assert "Slide 1" in result
        assert "Welcome to the presentation" in result

    def test_reads_table_from_slide(self, fake_engine, sample_pptx_meta):
        tool = ReadPptxSlideTextTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_pptx_meta.file_id, slide=2
        ))
        # Table should be rendered as markdown
        assert "Metric" in result
        assert "Revenue" in result
        assert "$1.2M" in result
        assert "|" in result  # markdown table pipes

    def test_slide_out_of_range(self, fake_engine, sample_pptx_meta):
        tool = ReadPptxSlideTextTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_pptx_meta.file_id, slide=99
        ))
        assert "Error" in result
        assert "out of range" in result

    def test_invalid_slide_number(self, fake_engine, sample_pptx_meta):
        tool = ReadPptxSlideTextTool(fake_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_pptx_meta.file_id, slide=0
        ))
        assert "Error" in result


# -----------------------------------------------------------------------------
# Registration
# -----------------------------------------------------------------------------


class _FakeManager:
    def __init__(self):
        self.tools = []

    def register_tool(self, tool):
        self.tools.append(tool)

    def register_function(self, **kwargs):
        pass


class TestRegistration:
    def test_excel_registers_both_tools(self, fake_engine):
        mgr = _FakeManager()
        ok = register_excel(mgr, fake_engine)
        assert ok is True
        names = {t.name for t in mgr.tools}
        assert "list_excel_sheets" in names
        assert "read_excel_sheet" in names

    def test_pptx_registers_both_tools(self, fake_engine):
        mgr = _FakeManager()
        ok = register_pptx(mgr, fake_engine)
        assert ok is True
        names = {t.name for t in mgr.tools}
        assert "list_pptx_slides" in names
        assert "read_pptx_slide_text" in names

    def test_excel_returns_false_without_engine(self):
        assert register_excel(_FakeManager(), None) is False

    def test_pptx_returns_false_without_engine(self):
        assert register_pptx(_FakeManager(), None) is False


# -----------------------------------------------------------------------------
# Workspace-path resolution (v1.18.7) — both Excel and PPTX tools must
# accept `path=` as an alternative to `file_id=` so they work on files
# landed via /files/upload or the file-tree drag-drop (no SessionFileStore
# entry exists for those).
# -----------------------------------------------------------------------------


@pytest.fixture
def workspace_engine(tmp_path, store):
    """Engine stub with a workspace dir AND a SessionFileStore."""
    wd = tmp_path / "workspace"
    wd.mkdir()
    return SimpleNamespace(file_store=store, get_working_dir=lambda: str(wd))


class TestWorkspacePathResolution:
    def test_list_excel_sheets_via_path(self, workspace_engine, tmp_path):
        xlsx = _make_xlsx({"Sheet1": [["a", "b"], [1, 2]]})
        (tmp_path / "workspace" / "data.xlsx").write_bytes(xlsx)
        tool = ListExcelSheetsTool(workspace_engine)
        result = asyncio.run(tool.execute(path="data.xlsx"))
        assert "data.xlsx" in result
        assert "Sheet1" in result
        assert "Error" not in result

    def test_read_excel_sheet_via_path(self, workspace_engine, tmp_path):
        xlsx = _make_xlsx({"Sheet1": [["a", "b"], [1, 2]]})
        (tmp_path / "workspace" / "data.xlsx").write_bytes(xlsx)
        tool = ReadExcelSheetTool(workspace_engine)
        result = asyncio.run(tool.execute(path="data.xlsx", sheet="Sheet1"))
        assert "a" in result and "b" in result
        assert "Error" not in result

    def test_list_pptx_slides_via_path(self, workspace_engine, tmp_path):
        pptx = _make_pptx([{"title": "Hello"}])
        (tmp_path / "workspace" / "deck.pptx").write_bytes(pptx)
        tool = ListPptxSlidesTool(workspace_engine)
        result = asyncio.run(tool.execute(path="deck.pptx"))
        assert "deck.pptx" in result
        assert "1 slide" in result
        assert "Error" not in result

    def test_read_pptx_slide_text_via_path(self, workspace_engine, tmp_path):
        pptx = _make_pptx([{"title": "Greetings", "content": "Body text here"}])
        (tmp_path / "workspace" / "deck.pptx").write_bytes(pptx)
        tool = ReadPptxSlideTextTool(workspace_engine)
        result = asyncio.run(tool.execute(path="deck.pptx", slide=1))
        assert "Greetings" in result or "Body text" in result

    def test_path_escape_rejected(self, workspace_engine, tmp_path):
        # File exists at tmp_path level (outside workspace) — must be rejected.
        pptx = _make_pptx([{"title": "Secret"}])
        (tmp_path / "outside.pptx").write_bytes(pptx)
        tool = ListPptxSlidesTool(workspace_engine)
        result = asyncio.run(tool.execute(path="../outside.pptx"))
        assert "Error" in result
        assert "outside the working directory" in result

    def test_both_file_id_and_path_rejected(self, workspace_engine, sample_pptx_meta, tmp_path):
        # Sending both at once must be an error, not pick-one silently.
        pptx = _make_pptx([{"title": "Hello"}])
        (tmp_path / "workspace" / "deck.pptx").write_bytes(pptx)
        tool = ListPptxSlidesTool(workspace_engine)
        result = asyncio.run(tool.execute(
            file_id=sample_pptx_meta.file_id, path="deck.pptx"
        ))
        assert "Error" in result
        assert "not both" in result.lower()

    def test_neither_arg_rejected(self, workspace_engine):
        tool = ListPptxSlidesTool(workspace_engine)
        result = asyncio.run(tool.execute())
        assert "Error" in result
        assert "file_id" in result or "path" in result


# -----------------------------------------------------------------------------
# render_pptx_slide: artifact-save behavior (v1.18.7)
# -----------------------------------------------------------------------------


class TestRenderPptxSlideArtifactFlow:
    """Rendered slide PNG must go into SessionFileStore, NOT inline as
    base64 — the inline approach cost ~100K tokens per call and blew
    the context window on every single render."""

    def _render(self, workspace_engine, tmp_path, slide_num=1):
        from ppxai.common.libreoffice import libreoffice_can_read
        from ppxai.engine.tools.builtin.pptx_tools import (
            RenderPptxSlideTool,
            _libreoffice_available,
        )
        if not _libreoffice_available():
            pytest.skip("LibreOffice not installed — render tests need it")
        # A snap-confined LibreOffice can't read tmp_path (outside $HOME), so a
        # convert silently yields no output (exit 0). Skip rather than report a
        # false regression — the coder image's apt libreoffice is unconfined and
        # this passes there (verified in-pod 2026-08-11).
        if not libreoffice_can_read(tmp_path):
            pytest.skip(
                "LibreOffice cannot read the test tmp dir (snap confinement) — "
                "render works with an unconfined install (e.g. the coder image)"
            )
        pptx_bytes = _make_pptx([{"title": "Hello", "content": "Body"}])
        (tmp_path / "workspace" / "deck.pptx").write_bytes(pptx_bytes)
        tool = RenderPptxSlideTool(workspace_engine)
        return asyncio.run(tool.execute(path="deck.pptx", slide=slide_num))

    def test_render_saves_to_session_store(self, workspace_engine, tmp_path):
        result = self._render(workspace_engine, tmp_path)
        # No inline base64 — that was the bug.
        assert "data:image/png;base64," not in result
        # Compact reference present.
        assert "file_id=" in result
        # PNG actually saved.
        store_meta = workspace_engine.file_store.list_all()
        pngs = [m for m in store_meta if m.media_type == "image/png"]
        assert len(pngs) >= 1
        assert any(m.name.endswith("_slide_1.png") for m in pngs)

    def test_render_result_is_compact(self, workspace_engine, tmp_path):
        # Result string must be small — ~1KB max, not 100K+.
        result = self._render(workspace_engine, tmp_path)
        # Allow ample headroom for description text, but bytes far
        # below the data-URI ceiling.
        assert len(result) < 2048, f"render result too large ({len(result)} bytes)"

    def test_render_falls_back_to_data_uri_without_store(self, tmp_path):
        # Engine stub WITHOUT file_store: legacy inline path stays
        # available so non-store callers still get pixels.
        from ppxai.common.libreoffice import libreoffice_can_read
        from ppxai.engine.tools.builtin.pptx_tools import (
            RenderPptxSlideTool,
            _libreoffice_available,
        )
        if not _libreoffice_available():
            pytest.skip("LibreOffice not installed")
        if not libreoffice_can_read(tmp_path):
            pytest.skip(
                "LibreOffice cannot read the test tmp dir (snap confinement)"
            )
        wd = tmp_path / "workspace"
        wd.mkdir()
        engine = SimpleNamespace(get_working_dir=lambda: str(wd))
        pptx_bytes = _make_pptx([{"title": "Hi"}])
        (wd / "deck.pptx").write_bytes(pptx_bytes)
        tool = RenderPptxSlideTool(engine)
        result = asyncio.run(tool.execute(path="deck.pptx", slide=1))
        # Without a store, the tool falls back to inline data URI.
        assert "data:image/png;base64," in result
